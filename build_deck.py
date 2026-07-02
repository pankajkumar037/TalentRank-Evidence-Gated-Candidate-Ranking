#!/usr/bin/env python3
"""
build_deck.py — Populate IdeaSubmissionTemplate.pptx with the TalentRank solution.

Preserves every template background (branding intact) and fills the white content area
of each slide with professionally formatted content, plus builds diagrams for the
Workflow, Architecture, and Results slides. Output: TalentRank_IdeaSubmission.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEMPLATE = "IdeaSubmissionTemplate.pptx"
OUT = "TalentRank_IdeaSubmission.pptx"

INK      = RGBColor(0x20, 0x27, 0x29)
MUTED    = RGBColor(0x5A, 0x62, 0x70)
PURPLE   = RGBColor(0x6A, 0x2C, 0xE0)
BLUE     = RGBColor(0x2A, 0x39, 0xD6)
GOOD     = RGBColor(0x12, 0x8A, 0x5E)
TILE_BG  = RGBColor(0xF3, 0xF0, 0xFD)
TILE_BG2 = RGBColor(0xEE, 0xF1, 0xFE)
TILE_BRD = RGBColor(0xD9, 0xCC, 0xF7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FONT     = "Manrope"


# --------------------------------------------------------------------------
def _run(p, text, *, size=12.5, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run(); r.text = text
    f = r.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return r


def _spacing(p, before=0.0, after=5.0, line=1.06):
    p.space_before = Pt(before); p.space_after = Pt(after); p.line_spacing = line


def _no_bullet(p):
    """Suppress any inherited list bullet so only our manual glyphs show."""
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _set_line(shape, label, value):
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.clear(); _spacing(p, after=0, line=1.02)
    _run(p, label, size=13, color=INK, bold=True)
    _run(p, value, size=13, color=INK, bold=False)


def body_box(slide):
    best = None
    for sh in slide.shapes:
        if sh.has_text_frame and Emu(sh.top).inches > 1.2:
            best = sh
    return best


def clear_tf(tf):
    p0 = tf.paragraphs[0]
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    p0.clear()
    return tf


def _style_run(p, text, style, size=12.5, lead=False):
    if style == "a":
        _run(p, text, size=size, color=PURPLE, bold=True)
    elif style == "blue":
        _run(p, text, size=size, color=BLUE, bold=True)
    elif style == "good":
        _run(p, text, size=size, color=GOOD, bold=True)
    elif style == "b":
        _run(p, text, size=size, color=INK, bold=True)
    elif style == "mut":
        _run(p, text, size=size, color=MUTED, bold=False)
    else:
        _run(p, text, size=size, color=(PURPLE if lead else INK), bold=lead)


def add_bullets(slide, blocks, base=12.5):
    box = body_box(slide)
    tf = box.text_frame; tf.word_wrap = True
    clear_tf(tf)
    first = True
    for blk in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _no_bullet(p)
        if "lead" in blk:
            _spacing(p, before=0, after=7, line=1.03)
            for text, style in blk["lead"]:
                _style_run(p, text, style, size=base + 1.5, lead=True)
        elif "gap" in blk:
            _spacing(p, before=0, after=0, line=0.5)
            _run(p, " ", size=5)
        else:
            lvl = blk.get("lvl", 0)
            _spacing(p, before=0, after=(5 if lvl == 0 else 3), line=1.03)
            if lvl == 0:
                _run(p, "▪  ", size=base, color=PURPLE, bold=True); sz = base
            else:
                _run(p, "        –  ", size=base - 1, color=MUTED); sz = base - 1
            for text, style in blk["b"]:
                _style_run(p, text, style, size=sz)


def tile(slide, x, y, w, h, big, label, *, big_color=PURPLE, bg=TILE_BG, sub=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = bg
    sp.line.color.rgb = TILE_BRD; sp.line.width = Pt(1)
    sp.adjustments[0] = 0.12; sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    for m in ("left", "right"): setattr(tf, f"margin_{m}", Inches(0.06))
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _spacing(p, after=0, line=0.9)
    _run(p, big, size=20, color=big_color, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; _spacing(p2, before=1, after=0, line=0.9)
    _run(p2, label, size=9, color=INK, bold=True)
    if sub:
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; _spacing(p3, before=0, after=0, line=0.85)
        _run(p3, sub, size=7.5, color=MUTED)
    return sp


def flow_box(slide, x, y, w, h, lines, *, bg=TILE_BG, brd=TILE_BRD):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = bg
    sp.line.color.rgb = brd; sp.line.width = Pt(1.25)
    sp.adjustments[0] = 0.16; sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, sz, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; _spacing(p, after=0, line=0.92)
        _run(p, txt, size=sz, color=col, bold=bold)
    return sp


def arrow(slide, x1, y1, x2, y2, color=PURPLE, w=2.25):
    cn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


# ==========================================================================
prs = Presentation(TEMPLATE)
S = list(prs.slides)

# ---- SLIDE 1 — Title ------------------------------------------------------
s1 = S[0]
for sh in s1.shapes:
    if sh.has_text_frame:
        t = sh.text_frame.paragraphs[0].text
        if t.startswith("Team Name"):
            _set_line(sh, "Team Name :  ", "TalentRank")
        elif t.startswith("Team Leader"):
            _set_line(sh, "Team Leader Name :  ", "Priyanshu Saraswat")
        elif t.startswith("Problem Statement"):
            _set_line(sh, "Problem Statement :  ",
                      "Intelligent Candidate Discovery & Ranking — rank the top 100 of "
                      "100,000 candidates for a Senior AI Engineer (Founding Team) role, "
                      "under a CPU-only, no-network, 5-minute budget.")
mb = s1.shapes.add_textbox(Inches(0.34), Inches(5.02), Inches(9.3), Inches(0.4))
tf = mb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; _spacing(p, after=0, line=1.0)
_run(p, "Team : ", size=11, color=INK, bold=True)
_run(p, "Priyanshu Saraswat (Leader) · Sujal Awasthi · Pankaj Kumar", size=11, color=INK)

# ---- SLIDE 2 — Solution Overview -----------------------------------------
add_bullets(S[1], [
    {"lead": [("An evidence-gated ranking engine — credit attaches to ", ""),
              ("described work", "a"), (", never to claimed labels.", "")]},
    {"b": [("Scoring: ", "b"), ("final = gate × quality × availability", "a"),
           ("   where   quality = 0.55·evidence + 0.20·depth + 0.25·fit", "")]},
    {"gap": 1},
    {"b": [("Reads career-history prose", "b"),
           (" — plain language scores like buzzwords: \"connect users to the most relevant "
            "matches\" == \"RAG / Pinecone\".", "")]},
    {"b": [("Ignores 63 noise skill tags", "b"),
           (" (each ~12,000×); trusts only ", ""), ("14 rare curated tags", "a"),
           (" (<10 uses each) as corroboration.", "")]},
    {"b": [("Removes impossible \"honeypot\" profiles", "b"),
           (" and down-weights ghosts via behavioral signals.", "")]},
    {"b": [("Explainable, deterministic, CPU-only, no LLM at rank time", "b"),
           (" — can't hallucinate, and beats the traps the challenge plants.", "")]},
])

# ---- SLIDE 3 — JD Understanding & Candidate Evaluation --------------------
add_bullets(S[2], [
    {"lead": [("We rank on what the JD ", ""), ("means", "a"), (", not the keywords it lists.", "")]},
    {"b": [("Four must-haves extracted: ", "b"),
           ("(1) shipped ranking/recommendation  (2) production embeddings/retrieval  "
            "(3) vector-DB / hybrid search  (4) rigorous evaluation (NDCG · MRR · MAP · A/B).", "")]},
    {"b": [("Hard disqualifiers (JD-stated): ", "b"),
           ("consultancy-only career · pure academia w/o production · abroad & won't "
            "relocate · LLM-wrapper-only · wrong domain (CV/speech/robotics w/o NLP).", "")]},
    {"b": [("Signal priority: ", "b"),
           ("described work (primary) → behavioral availability (login recency, recruiter "
            "response, open-to-work, notice, interview completion) → location (Noida/Pune), "
            "seniority, ~7-yr experience curve.", "")]},
    {"b": [("Beyond keywords: ", "b"),
           ("a plain-language Tier-5 who \"built a recommendation system at a product "
            "company\" outranks a Marketing Manager carrying every AI keyword.", "")]},
])

# ---- SLIDE 4 — Ranking Methodology ----------------------------------------
add_bullets(S[3], [
    {"lead": [("Retrieve → score → modulate → order", "a")]},
    {"b": [("Retrieve: ", "b"),
           ("stream 100k JSONL → C-speed substring prefilter → regex must-have detection "
            "on description prose, recency-weighted per role.", "")]},
    {"b": [("Score each candidate:", "b")]},
    {"b": [("evidence", "a"), (" — proof of the 4 must-haves  ", ""),
           ("(0.55)", "mut")], "lvl": 1},
    {"b": [("depth", "a"), (" — non-saturating differentiator: scale · breadth · ownership · "
            "technique specificity · recency  ", ""), ("(0.20) — the NDCG@10 lever", "mut")], "lvl": 1},
    {"b": [("fit", "a"), (" — experience curve · ML-years @ product · hands-on · location · "
            "seniority · JD penalties  ", ""), ("(0.25)", "mut")], "lvl": 1},
    {"b": [("Modulate: ", "b"),
           ("× gate (hard deal-breakers → ≈0)   × availability (behavioral → ×[0.75, 1.0]).", "")]},
    {"b": [("Order: ", "b"),
           ("multiplicative compose → sort by score → deterministic tie-break by "
            "candidate_id → top 100.", "")]},
])

# ---- SLIDE 5 — Explainability & Data Validation ---------------------------
add_bullets(S[4], [
    {"lead": [("Every rank is explained from facts — never generated free-text.", "a")]},
    {"b": [("Grounded reasoning", "b"),
           (" per candidate: templated from the exact signals the scorer used + a verbatim "
            "quote from the candidate's own text → ", ""), ("100/100 distinct.", "good")]},
    {"b": [("Anti-hallucination by construction: ", "b"),
           ("every claim maps to profile data; tone matches rank; honest concerns surfaced "
            "(notice period, ghost, management-heavy).", "")]},
    {"b": [("Honeypots", "b"),
           (" (impossible profiles: \"expert\" skill w/ 0 months, tenure > career span, "
            "years ≫ history) → forced to 0 and removed → ", ""),
           ("0 in the top 100.", "good")]},
    {"b": [("Trap resistance: ", "b"),
           ("keyword-stuffers gated, 63 noise tags ignored → skill-tag trap beaten ", ""),
           ("2.4×.", "good")]},
])

# ---- SLIDE 9 — Technologies Used ------------------------------------------
add_bullets(S[8], [
    {"lead": [("Chosen for the CPU-only / no-network / 5-min constraint.", "a")]},
    {"b": [("Core: ", "b"), ("Python 3.11 · NumPy · orjson", "a"),
           (" (fast streaming of the 465 MB JSONL).", "")]},
    {"b": [("Ranking engine: ", "b"),
           ("hand-built deterministic rules — regex must-have detectors + recency weighting. "
            "No ML inference, no network at rank time.", "")]},
    {"b": [("Output & QA: ", "b"), ("openpyxl", "a"), (" (.xlsx) · csv (stdlib) · ", ""),
           ("pytest", "a"), (" (14 adversarial regression tests).", "")]},
    {"b": [("Sandbox: ", "b"), ("Streamlit", "a"),
           (" — upload ≤100 profiles, live gate × quality × availability breakdown.", "")]},
    {"b": [("Built, measured & deliberately excluded: ", "b"),
           ("sentence-transformers / all-MiniLM-L6-v2", "mut"),
           (" — offline embeddings recreated keyword bias, so we dropped it.", "")]},
])

# ---- SLIDE 10 — Submission Assets -----------------------------------------
add_bullets(S[9], [
    {"lead": [("Everything needed to reproduce and verify, end to end.", "a")]},
    {"b": [("GitHub repo: ", "b"),
           ("full source · README with one-command reproduce · requirements.txt · "
            "submission_metadata.yaml.   ", ""), ("[ add repo link ]", "blue")]},
    {"b": [("Deliverables: ", "b"), ("submission.csv", "a"),
           (" (validator-passed)  +  ", ""), ("submission.xlsx", "a"),
           (" (portal upload).", "")]},
    {"b": [("Reproduce: ", "b"),
           ("python rank.py --candidates ./data/candidates.jsonl --out ./submission.csv", "a")]},
    {"b": [("Validation: ", "b"),
           ("validate_submission.py · verify_submission.py (0 honeypots) · evaluate.py · "
            "pytest (14/14).", "")]},
    {"b": [("Sandbox / demo: ", "b"), ("Streamlit app  ", ""),
           ("[ add hosted link ]", "blue")]},
])

# ---- SLIDE 6 — End-to-End Workflow (diagram) ------------------------------
s6 = S[5]
bb = body_box(s6)
if bb is not None:
    clear_tf(bb.text_frame)
row = [
    ("JD + signals", "parsed to\nintent"),
    ("JD-cited config", "weights &\nthresholds"),
    ("Stream 100k", "candidate\nprofiles"),
    ("Honeypot + gate", "fast-path\nfilter"),
    ("evidence · depth · fit", "quality\nscore"),
    ("× availability", "behavioral\nmultiplier"),
    ("Rank + reason", "top-100\n+ CSV/XLSX"),
]
n = len(row); gap = Inches(0.12)
total_w = Inches(9.1)
bw = Emu(int((total_w - gap * (n - 1)) / n))
x = Inches(0.45); yb = Inches(2.55); bh = Inches(1.15)
cy = yb + Emu(int(bh / 2))
pastels = [TILE_BG, TILE_BG, TILE_BG2, TILE_BG, TILE_BG, TILE_BG2, TILE_BG]
for i, (title, sub) in enumerate(row):
    bx = Emu(x + (bw + gap) * i)
    flow_box(s6, bx, yb, bw, bh,
             [(title, 10, INK, True), (sub, 8, MUTED, False)],
             bg=pastels[i])
    if i < n - 1:
        ax1 = Emu(bx + bw); ax2 = Emu(bx + bw + gap)
        arrow(s6, ax1, cy, ax2, cy, color=PURPLE, w=2)
cap = s6.shapes.add_textbox(Inches(0.45), Inches(4.05), Inches(9.1), Inches(0.5))
cp = cap.text_frame.paragraphs[0]; _spacing(cp, after=0, line=1.0)
_run(cp, "One deterministic pass — CPU-only, no network, ~93 s on 100k. Pre-computation "
         "(if any) happens offline; the ranking step alone produces the submission.",
     size=10, color=MUTED, italic=True)

# ---- SLIDE 7 — System Architecture (diagram) ------------------------------
s7 = S[6]
# config band (feeds all)
flow_box(s7, Inches(0.5), Inches(1.55), Inches(9.0), Inches(0.5),
         [("src/jd_config.py  —  every weight & threshold, each citing the JD sentence that justifies it",
           10, PURPLE, True)], bg=TILE_BG, brd=TILE_BRD)
# loader
flow_box(s7, Inches(0.5), Inches(2.25), Inches(1.7), Inches(0.95),
         [("loader.py", 10.5, INK, True), ("stream JSONL\n→ clean facts", 8, MUTED, False)], bg=TILE_BG2)
# two analysis columns
gates_lines = [("gates · honeypot", 9.5, INK, True), ("availability", 9.5, INK, True),
               ("deal-breakers,\nimpossible profiles,\nbehavioral signals", 8, MUTED, False)]
ev_lines = [("evidence (+curated_tags)", 9.5, INK, True), ("fit", 9.5, INK, True),
            ("4 must-haves, depth,\nexperience · location ·\nseniority · penalties", 8, MUTED, False)]
flow_box(s7, Inches(2.55), Inches(2.15), Inches(2.35), Inches(1.15), gates_lines, bg=TILE_BG)
flow_box(s7, Inches(2.55), Inches(3.35), Inches(2.35), Inches(0.0)+Inches(0.9),
         [("", 1, INK, False)], bg=WHITE, brd=WHITE)  # spacer (invisible)
flow_box(s7, Inches(5.05), Inches(2.15), Inches(2.4), Inches(1.15), ev_lines, bg=TILE_BG)
# scorer
flow_box(s7, Inches(7.65), Inches(2.25), Inches(1.85), Inches(0.95),
         [("scorer.py", 10.5, PURPLE, True), ("gate × quality\n× availability", 8.5, INK, True)], bg=TILE_BG)
# outputs row
flow_box(s7, Inches(2.55), Inches(3.62), Inches(2.35), Inches(0.75),
         [("reasoning.py", 9.5, INK, True), ("grounded, varied\nexplanations", 8, MUTED, False)], bg=TILE_BG2)
flow_box(s7, Inches(5.05), Inches(3.62), Inches(2.4), Inches(0.75),
         [("outputs", 9.5, INK, True), ("submission.csv + .xlsx\n(validated)", 8, MUTED, False)], bg=TILE_BG2)
flow_box(s7, Inches(7.65), Inches(3.62), Inches(1.85), Inches(0.75),
         [("Streamlit sandbox", 9, INK, True), ("+ adversarial tests", 8, MUTED, False)], bg=TILE_BG2)
# arrows
arrow(s7, Inches(2.2), Inches(2.72), Inches(2.55), Inches(2.72), color=PURPLE, w=2)
arrow(s7, Inches(4.9), Inches(2.72), Inches(5.05), Inches(2.72), color=PURPLE, w=2)
arrow(s7, Inches(7.45), Inches(2.72), Inches(7.65), Inches(2.72), color=PURPLE, w=2)
arrow(s7, Inches(8.55), Inches(3.20), Inches(8.55), Inches(3.62), color=BLUE, w=2)
cap7 = s7.shapes.add_textbox(Inches(0.5), Inches(4.55), Inches(9.0), Inches(0.5))
cp7 = cap7.text_frame.paragraphs[0]; _spacing(cp7, after=0, line=1.0)
_run(cp7, "Modular, testable, and offline. No hosted LLM, no vector service, no network — "
          "the whole ranker is pure Python rules over precomputed facts.",
     size=10, color=MUTED, italic=True)

# ---- SLIDE 8 — Results & Performance (tiles) ------------------------------
s8 = S[7]
bb8 = body_box(s8)
if bb8 is not None:
    clear_tf(bb8.text_frame)
# quality metric tiles (independent silver-label judge)
lbl = s8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.3))
lp = lbl.text_frame.paragraphs[0]; _spacing(lp, after=0, line=1.0)
_run(lp, "Estimated quality — independent silver-label judge (kept blind to our curated tags):",
     size=10.5, color=INK, bold=True)
tw, th, gy = Inches(1.66), Inches(0.95), Inches(0.16)
xs = Inches(0.5); yq = Inches(1.86)
tiles = [("0.904", "Composite", PURPLE, "0.50·N@10 + 0.30·N@50 + …"),
         ("0.862", "NDCG@10", BLUE, "top-10 quality"),
         ("0.909", "NDCG@50", BLUE, "top-50 quality"),
         ("1.000", "MAP", GOOD, "all relevance levels"),
         ("1.000", "P@10", GOOD, "top-10 relevant")]
for i, (b, l, c, sub) in enumerate(tiles):
    tile(s8, Emu(xs + (tw + gy) * i), yq, tw, th, b, l, big_color=c,
         bg=(TILE_BG if c != GOOD else RGBColor(0xEC,0xF7,0xF1)), sub=sub)
# baseline comparison line
bl = s8.shapes.add_textbox(Inches(0.5), Inches(2.95), Inches(9.0), Inches(0.35))
blp = bl.text_frame.paragraphs[0]; _spacing(blp, after=0, line=1.0)
_run(blp, "vs baselines — ", size=10.5, color=INK, bold=True)
_run(blp, "keyword matcher 0.726", size=10.5, color=MUTED)
_run(blp, "   ·   ", size=10.5, color=MUTED)
_run(blp, "skill-tag trap 0.383", size=10.5, color=MUTED)
_run(blp, "   ·   ", size=10.5, color=MUTED)
_run(blp, "random 0.012", size=10.5, color=MUTED)
# integrity + compute tiles row
lbl2 = s8.shapes.add_textbox(Inches(0.5), Inches(3.42), Inches(9.0), Inches(0.3))
lp2 = lbl2.text_frame.paragraphs[0]; _spacing(lp2, after=0, line=1.0)
_run(lp2, "Integrity & compute — comfortably inside every challenge limit:",
     size=10.5, color=INK, bold=True)
tiles2 = [("0", "honeypots in top 100", GOOD, "DQ if >10%"),
          ("8 / 8", "apex golds surfaced", PURPLE, "was 3/8"),
          ("14/14", "adversarial tests", GOOD, "all pass"),
          ("93 s", "runtime / 300 s", BLUE, "CPU-only"),
          ("1.35 GB", "RAM / 16 GB", BLUE, "deterministic")]
yc = Inches(3.74)
for i, (b, l, c, sub) in enumerate(tiles2):
    tile(s8, Emu(xs + (tw + gy) * i), yc, tw, th, b, l, big_color=c,
         bg=(RGBColor(0xEC,0xF7,0xF1) if c == GOOD else TILE_BG2), sub=sub)
cap8 = s8.shapes.add_textbox(Inches(0.5), Inches(4.82), Inches(9.0), Inches(0.35))
cp8 = cap8.text_frame.paragraphs[0]; _spacing(cp8, after=0, line=1.0)
_run(cp8, "Byte-identical output across runs. The judge confirms our top-10 are tier 4–5 "
          "on prose alone — an independent check, not a circular one.",
     size=9.5, color=MUTED, italic=True)

# ---- SLIDE 11 — Closing (team footer) -------------------------------------
s11 = S[10]
fb = s11.shapes.add_textbox(Inches(0.5), Inches(5.02), Inches(9.0), Inches(0.5))
fp = fb.text_frame.paragraphs[0]; fp.alignment = PP_ALIGN.CENTER; _spacing(fp, after=0, line=1.0)
_run(fp, "TalentRank", size=11, color=WHITE, bold=True)
_run(fp, "   ·   Priyanshu Saraswat (Leader) · Sujal Awasthi · Pankaj Kumar",
     size=10, color=WHITE)

prs.save(OUT)
print(f"[done] wrote {OUT} with {len(S)} slides")
